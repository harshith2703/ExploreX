import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Configure 16:9 Widescreen slides (13.33 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors Palette
    # Dark Mode (Title & Conclusion slides)
    dark_bg = RGBColor(11, 23, 27)         # Deep Teal-Black #0b171b
    dark_card_bg = RGBColor(20, 39, 44)    # #14272c
    dark_card_border = RGBColor(38, 70, 78) # #26464e
    
    # Light Mode (Content slides)
    light_bg = RGBColor(248, 250, 252)     # Off-white #f8fafc
    light_card_bg = RGBColor(255, 255, 255) # White #ffffff
    light_card_border = RGBColor(226, 232, 240) # Slate #e2e8f0
    
    # Text and Accents
    primary_teal = RGBColor(11, 76, 82)    # Deep Teal #0b4c52
    mint_green = RGBColor(13, 148, 136)   # Mint #0d9488
    gold_accent = RGBColor(217, 119, 6)    # Amber Gold #d97706
    white_text = RGBColor(255, 255, 255)   # White
    slate_dark = RGBColor(51, 65, 85)      # Slate Dark Text #334155
    slate_light = RGBColor(100, 116, 139)  # Slate Light Body #64748b
    
    blank_layout = prs.slide_layouts[6]
    
    # Image Paths
    screenshots_dir = r"c:\Users\harsh\OneDrive\Desktop\ExploreX\documentation\screenshots"
    home_img_path = os.path.join(screenshots_dir, "home.png")
    login_img_path = os.path.join(screenshots_dir, "login.png")
    register_img_path = os.path.join(screenshots_dir, "register.png")
    categories_img_path = os.path.join(screenshots_dir, "categories.png")
    
    def set_solid_background(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, title_text, is_dark=False):
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = white_text if is_dark else primary_teal

    def add_card(slide, left, top, width, height, title, body, is_dark_slide=False, accent_border=False):
        # Draw card container
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        
        # Color based on slide background type
        if is_dark_slide:
            shape.fill.fore_color.rgb = dark_card_bg
            shape.line.color.rgb = mint_green if accent_border else dark_card_border
        else:
            shape.fill.fore_color.rgb = light_card_bg
            shape.line.color.rgb = mint_green if accent_border else light_card_border
            
        shape.line.width = Pt(1.5)
        
        # Draw card text
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        # Add Card Title
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = mint_green if is_dark_slide else primary_teal
        p_title.space_after = Pt(8)
        
        # Add Card Body
        if body:
            p_body = tf.add_paragraph()
            p_body.text = body
            p_body.font.name = "Segoe UI"
            p_body.font.size = Pt(11)
            p_body.font.color.rgb = white_text if is_dark_slide else slate_dark
            
        return shape

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide1, dark_bg)
    
    # Project Title
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ExploreX"
    p.font.name = "Segoe UI"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = mint_green
    
    p_sub = tf.add_paragraph()
    p_sub.text = "A Tourist Information System"
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(28)
    p_sub.font.color.rgb = white_text
    p_sub.space_after = Pt(14)
    
    # Subtitle / Pitch
    p2 = tf.add_paragraph()
    p2.text = "A production-ready full-stack travel portal built with React.js, Spring Boot, and MySQL — enabling tourist destination discovery, dynamic maps lookup, and secure guided tour package bookings."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = slate_light
    p2.space_before = Pt(8)
    
    # Badges / Chips
    add_card(slide1, 1.0, 5.8, 2.5, 0.45, "FULL-STACK PROJECT", "", is_dark_slide=True)
    add_card(slide1, 3.7, 5.8, 3.4, 0.45, "SPRING SECURITY + JWT + REACT", "", is_dark_slide=True)
    
    # ==========================================
    # SLIDE 2: Introduction & Objectives (Light Theme, Home Screenshot)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide2, light_bg)
    add_slide_header(slide2, "Introduction & Objectives")
    
    # Description Text
    txBox = slide2.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.4), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ExploreX solves travel planning friction by offering tourists a unified platform to explore curated places, review trip details, and request bookings directly. It features separate, secure portal views for travelers and administrative managers."
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = slate_dark
    
    # Left Column: Objectives Cards
    add_card(slide2, 1.0, 2.9, 5.4, 1.8, "Objectives", 
             "• Empower travelers with dynamic location discovery, ratings, and filters.\n"
             "• Simplify tour bookings through request forms and guides assignment.\n"
             "• Streamline portal management via responsive CRUD dashboards.")
    
    add_card(slide2, 1.0, 4.9, 5.4, 1.7, "Value Proposition", 
             "• Verified Reviews Engine: Aggregates tourist scores automatically.\n"
             "• High Security: Access control for users and admin panels via JWT.\n"
             "• Unified Catalog: Consolidated Beach, Mountain, Wildlife, and Heritage places.")
    
    # Right Column: Screenshot
    if os.path.exists(home_img_path):
        slide2.shapes.add_picture(home_img_path, Inches(6.9), Inches(1.6), width=Inches(5.4), height=Inches(5.0))
        # Draw a subtle border around the picture
        border = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.4), Inches(5.0))
        border.fill.background()
        border.line.color.rgb = light_card_border
        border.line.width = Pt(1.5)
        
    # ==========================================
    # SLIDE 3: Technologies Used (Light Theme)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide3, light_bg)
    add_slide_header(slide3, "Technologies Used")
    
    # Left descriptive card
    add_card(slide3, 1.0, 1.8, 3.6, 4.8, "ExploreX Stack", 
             "A robust, production-quality architecture utilizing enterprise frameworks on the backend and modular reactivity on the frontend to ensure performance, reliability, and responsive scaling.", 
             accent_border=True)
    
    # 4 Tech Category Cards
    add_card(slide3, 5.0, 1.8, 3.6, 2.2, "Frontend UI", "React.js · Vite · HTML5 · CSS3\nAxios Client · Context API\nBootstrap & Custom Theme Styles")
    add_card(slide3, 8.9, 1.8, 3.6, 2.2, "Backend Engine", "Spring Boot (Java) · Maven\nSpring Security · JWT Token Auth\nSpring Data JPA Hibernate")
    add_card(slide3, 5.0, 4.4, 3.6, 2.2, "Database", "MySQL Server Relational DB\n`schema.sql` DDL definitions\n`data.sql` Seed records & Roles")
    add_card(slide3, 8.9, 4.4, 3.6, 2.2, "Tools & APIs", "RESTful Architecture · Git/GitHub\nPostman API Test Suite\nVS Code · Spring Tool Suite")
    
    # ==========================================
    # SLIDE 4: System Architecture (Light Theme)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide4, light_bg)
    add_slide_header(slide4, "System Architecture")
    
    # Subtitle
    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Multilayer relational data flow showing front-to-back integration, security filtering, and persistence mapping."
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = slate_light
    
    # Flow Nodes
    add_card(slide4, 5.2, 3.6, 2.8, 1.4, "System Layers\nData Flow Map", "", accent_border=True)
    
    add_card(slide4, 5.2, 2.1, 2.8, 1.2, "React Frontend (UI)", "Captures user input and renders responsive views.")
    add_card(slide4, 8.4, 3.7, 2.8, 1.2, "Spring Security Filter", "Authenticates endpoints using stateless JWT interceptors.")
    add_card(slide4, 8.4, 5.2, 2.8, 1.2, "Controller & Service", "Exposes REST endpoints & handles transaction logic.")
    add_card(slide4, 5.2, 5.2, 2.8, 1.2, "Data JPA Repositories", "Query translation layers to map Entities safely.")
    add_card(slide4, 2.0, 3.7, 2.8, 1.2, "MySQL Database", "Stores persistent relational schemas and tables.")
    add_card(slide4, 2.0, 2.1, 2.8, 1.2, "Axios / DTO API Calls", "Passes JSON requests and payloads over HTTP.")
    
    # Bottom text
    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡 Decoupled frontend-backend separation facilitates independent maintenance, horizontal scaling, and secure data pipelines."
    p.font.name = "Segoe UI"
    p.font.size = Pt(12)
    p.font.color.rgb = slate_light
    
    # ==========================================
    # SLIDE 5: Project Structure (Light Theme)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide5, light_bg)
    add_slide_header(slide5, "Project Structure")
    
    # Subheader
    txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Folder Organization"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = primary_teal
    
    p2 = tf.add_paragraph()
    p2.text = "The codebase is organized into four main directories to maintain clean separation of concerns."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(14)
    p2.font.color.rgb = slate_dark
    p2.space_before = Pt(6)
    
    # 4 Columns Cards
    add_card(slide5, 1.0, 3.0, 2.6, 3.5, "📁 frontend/", "• context/ (JWT auth context)\n• services/ (Axios calls)\n• pages/ (Home, Catalog, Details)\n• layouts/ (Nav & footers)\n• index.css (HSL themes)")
    add_card(slide5, 3.9, 3.0, 2.6, 3.5, "📁 backend/", "• controller/ (Endpoints)\n• service/ (Interface / Impl)\n• repository/ (JPA Layer)\n• entity/ (JPA Mapping)\n• security/jwt/ (JWT filter)\n• dto/ & mapper/ (DTO contracts)")
    add_card(slide5, 6.8, 3.0, 2.6, 3.5, "📁 database/", "• schema.sql\n(Tables DDL schema definitions)\n\n• data.sql\n(Roles, Admin/User seeds, default catalog values)")
    add_card(slide5, 9.7, 3.0, 2.6, 3.5, "📁 documentation/", "• INSTALLATION.md\n(Setup guidelines)\n\n• API_DOCUMENTATION.md\n(List of REST endpoints)\n\n• DIAGRAMS.md (ER maps)")

    # ==========================================
    # SLIDE 6: Database Design (Light Theme)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide6, light_bg)
    add_slide_header(slide6, "Database Design")
    
    # Title
    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Relational Schemas"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = primary_teal
    
    p2 = tf.add_paragraph()
    p2.text = "Normalized tables linked via relational constraints to maintain data integrity and referential constraints."
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(14)
    p2.font.color.rgb = slate_dark
    p2.space_before = Pt(6)
    
    # Left core cards
    add_card(slide6, 1.0, 3.0, 4.2, 1.6, "Users & Roles Tables", "• users: id, username, email, password, role_id\n• roles: id, name (ROLE_USER, ROLE_ADMIN)")
    add_card(slide6, 1.0, 4.8, 4.2, 1.6, "Places & Categories Tables", "• places: id, name, location, latitude, longitude, description, category_id\n• categories: id, name")
    
    # Info badge
    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(4.2), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "★ Star Ratings & bookings are dynamically re-aggregated on review submissions."
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = gold_accent
    
    # Right ERD visual representation container
    add_card(slide6, 6.0, 1.8, 6.3, 5.0, "Database Relationship Map", "", accent_border=True)
    
    # Users box
    add_card(slide6, 7.5, 2.3, 3.3, 1.5, "Users Table", "id (PK)\nusername\nemail\npassword")
    
    # Connector text
    txBox = slide6.shapes.add_textbox(Inches(7.5), Inches(3.9), Inches(3.3), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "↓↓   1 user has many bookings & reviews   ↓↓"
    p.font.name = "Segoe UI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = primary_teal
    p.alignment = 1 # Center
    
    # Bookings & Reviews boxes side-by-side
    add_card(slide6, 6.3, 4.4, 2.7, 1.7, "Bookings Table", "id (PK)\nbooking_date\nuser_id (FK)\ntourist_place_id (FK)\nguide_requested")
    add_card(slide6, 9.3, 4.4, 2.7, 1.7, "Reviews Table", "id (PK)\nstars_rating\ncomment_text\nuser_id (FK)\ntourist_place_id (FK)")
    
    # ==========================================
    # SLIDE 7: Frontend–Backend Connection (Light Theme)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide7, light_bg)
    add_slide_header(slide7, "Frontend–Backend Connection")
    
    # Subheading
    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "React communicates with Spring Boot via REST APIs, exchanging JSON data. Axios uses a request interceptor to attach JWT headers automatically."
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = slate_dark
    
    # Horizontal flow blocks
    add_card(slide7, 1.0, 2.2, 2.5, 0.7, "React Client (Vite)", "", accent_border=True)
    add_card(slide7, 3.8, 2.2, 2.5, 0.7, "Axios Interceptor", "")
    add_card(slide7, 6.6, 2.2, 2.5, 0.7, "Spring Security Filter", "")
    add_card(slide7, 9.4, 2.9, 2.9, 0.7, "Database Repository", "")
    
    # Rest endpoints table outline
    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(3.2), Inches(11.3), Inches(3.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Key REST API Endpoints"
    p.font.name = "Segoe UI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_teal
    p.space_after = Pt(10)
    
    # Drawing list as table representation
    rows = [
        ("POST", "/api/auth/register", "Register new tourist account"),
        ("POST", "/api/auth/login", "Authenticate credentials & return JWT token"),
        ("GET", "/api/places", "Fetch destinations list (supports filter & sorting params)"),
        ("POST", "/api/bookings", "Create a guided tour booking request (Auth needed)"),
        ("POST", "/api/places/{id}/reviews", "Post a traveler review & update ratings average"),
        ("POST", "/api/admin/places", "Create a new destination listing (Admin Only)")
    ]
    
    for method, endpoint, purpose in rows:
        p_row = tf.add_paragraph()
        p_row.text = f"•  [{method}]   {endpoint:<38}  →   {purpose}"
        p_row.font.name = "Courier New"
        p_row.font.size = Pt(12)
        p_row.font.bold = True
        p_row.font.color.rgb = slate_dark
        p_row.space_after = Pt(6)

    # ==========================================
    # SLIDE 8: Key Features & Authentication (Light Theme, Login Screenshot)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide8, light_bg)
    add_slide_header(slide8, "Key Features & Authentication")
    
    # 4 Grid Features on the left
    add_card(slide8, 1.0, 1.6, 2.7, 2.2, "User Authentication", 
             "• Secure sign-up/login.\n• Stateless JWT validation.\n• Client side routing guard by roles.")
             
    add_card(slide8, 4.0, 1.6, 2.7, 2.2, "Interactive Search", 
             "• Text filtering by place/state.\n• Category sorting (Beaches, Mountains, Heritage).\n• Order by rating stars.")
             
    add_card(slide8, 1.0, 4.2, 2.7, 2.2, "Reviews Engine", 
             "• Star grades from 1 to 5.\n• Written text comments.\n• Automatic real-time score recalculation.")
             
    add_card(slide8, 4.0, 4.2, 2.7, 2.2, "Guided Bookings", 
             "• Date selectors.\n• Traveler count multipliers.\n• Admin dashboard approval & rejection workflow.")
    
    # Right side: Welcome Back login screenshot
    if os.path.exists(login_img_path):
        slide8.shapes.add_picture(login_img_path, Inches(7.2), Inches(1.6), width=Inches(5.1), height=Inches(5.0))
        border = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(1.6), Inches(5.1), Inches(5.0))
        border.fill.background()
        border.line.color.rgb = light_card_border
        border.line.width = Pt(1.5)
        
    # ==========================================
    # SLIDE 9: Interactive Catalog & Bookings (Light Theme, Categories Screenshot)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide9, light_bg)
    add_slide_header(slide9, "Interactive Destination Catalog")
    
    # Left Description
    add_card(slide9, 1.0, 1.6, 5.4, 2.3, "Curated Categories", 
             "• Beaches: Coastal escapes with clean water activities.\n"
             "• Mountains: Breathtaking hill stations, mist-clad peaks, and tea gardens.\n"
             "• Heritage: Ancient palaces, historic forts, monuments, and archaeological sites.\n"
             "• Wildlife: Dense national parks, bio-diverse reserves, and safaris.")
             
    add_card(slide9, 1.0, 4.3, 5.4, 2.3, "Featured Places catalog", 
             "• Dynamic ratings system displays average score badges (e.g. ★ 4.8).\n"
             "• Visual cards show high-res galleries and quick descriptors.\n"
             "• Wishlist toggle lets users save locations to check later.")
    
    # Right side: Categories & Featured places screenshot
    if os.path.exists(categories_img_path):
        slide9.shapes.add_picture(categories_img_path, Inches(6.9), Inches(1.6), width=Inches(5.4), height=Inches(5.0))
        border = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.4), Inches(5.0))
        border.fill.background()
        border.line.color.rgb = light_card_border
        border.line.width = Pt(1.5)

    # ==========================================
    # SLIDE 10: Conclusion & Outcomes (Dark Theme, Register Screenshot)
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_solid_background(slide10, dark_bg)
    add_slide_header(slide10, "Conclusion & Outcomes", is_dark=True)
    
    # Left side: Text summary cards
    add_card(slide10, 1.0, 1.6, 5.4, 2.3, "Project Summary", 
             "ExploreX is a fully operational, responsive Tourist Information portal. It showcases a modern architecture leveraging Spring Boot REST security and dynamic React components to manage bookings, user reviews, and content dashboards securely.",
             is_dark_slide=True)
             
    add_card(slide10, 1.0, 4.3, 5.4, 2.3, "Key Achievements", 
             "• Role-based authorization using stateless JWT filters.\n"
             "• Automated averages calculation for stars reviews.\n"
             "• Seamless responsive UI matching across viewport sizes.",
             is_dark_slide=True, accent_border=True)
    
    # Right side: Sign up form screenshot
    if os.path.exists(register_img_path):
        slide10.shapes.add_picture(register_img_path, Inches(6.9), Inches(1.6), width=Inches(5.4), height=Inches(5.0))
        border = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.4), Inches(5.0))
        border.fill.background()
        border.line.color.rgb = dark_card_border
        border.line.width = Pt(1.5)
        
    # Save Presentation to Documents
    output_path = r"c:\Users\harsh\OneDrive\Documents\ExploreX_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    build_presentation()
